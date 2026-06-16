import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette (Premium Dark Theme)
    BG_COLOR = RGBColor(11, 19, 43)        # Slate 950 (Deep Blue/Gray)
    CARD_BG = RGBColor(28, 37, 65)         # Slate 800
    TEXT_MAIN = RGBColor(248, 250, 252)    # Off-white
    TEXT_MUTED = RGBColor(148, 163, 184)   # Cool gray
    ACCENT_CYAN = RGBColor(0, 180, 216)    # Vibrant light blue/cyan
    ACCENT_GREEN = RGBColor(74, 222, 128)   # Emerald green

    # Helper function to add a solid background to a slide
    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # Helper function to add standard slide titles
    def add_slide_header(slide, title, category="PARKFLOW AI"):
        # Category tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = "Outfit"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_CYAN
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        p_title = title_tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Outfit"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_MAIN

    # Helper to add standard bullets
    def add_textbox_bullets(slide, left, top, width, height, points):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for idx, (bold_txt, regular_txt) in enumerate(points):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(12)
            
            # Bold run
            r1 = p.add_run()
            r1.text = bold_txt
            r1.font.name = "Outfit"
            r1.font.bold = True
            r1.font.size = Pt(16)
            r1.font.color.rgb = TEXT_MAIN
            
            # Regular run
            r2 = p.add_run()
            r2.text = regular_txt
            r2.font.name = "Outfit"
            r2.font.size = Pt(16)
            r2.font.color.rgb = TEXT_MUTED
        return txBox

    # Helper to add card-like backgrounds for layout sections
    def add_card(slide, left, top, width, height):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = RGBColor(58, 80, 107)
        shape.line.width = Pt(1.5)
        return shape

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide1)

    # Title & Subtitle text box
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True

    # Category Tag
    p_tag = tf.paragraphs[0]
    p_tag.text = "FLIPKART GRIDLOCK HACKATHON 2.0"
    p_tag.font.name = "Outfit"
    p_tag.font.size = Pt(12)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT_CYAN
    p_tag.space_after = Pt(20)

    # Title
    p_title = tf.add_paragraph()
    p_title.text = "ParkFlow AI"
    p_title.font.name = "Outfit"
    p_title.font.size = Pt(54)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_MAIN
    p_title.space_after = Pt(10)

    # Subtitle
    p_sub = tf.add_paragraph()
    p_sub.text = "Intelligent Parking Congestion Management\nML Traffic Predictions meets Smart Dispatch Operations."
    p_sub.font.name = "Outfit"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.space_after = Pt(30)

    # Team Info
    p_team = tf.add_paragraph()
    p_team.text = "Developed by: Harshwardhan Bhaskar\nTeam: HB Technologies (Solo Participant)"
    p_team.font.name = "Outfit"
    p_team.font.size = Pt(14)
    p_team.font.bold = True
    p_team.font.color.rgb = ACCENT_GREEN

    # Image
    img1_path = r"docs/screenshots/smart_city_command_center.png"
    if os.path.exists(img1_path):
        slide1.shapes.add_picture(img1_path, Inches(7.6), Inches(1.5), Inches(5.0), Inches(4.5))

    # -------------------------------------------------------------
    # SLIDE 2: The Problem
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide2)
    add_slide_header(slide2, "The Gridlock Challenge: Reactive Enforcement")

    # Add bullets on left
    points2 = [
        ("Urban Congestion Bottlenecks: ", "Illegal parking in dense corridors block lanes, reducing road capacity by up to 65% and creating severe cascading traffic gridlocks."),
        ("Blind Random Patrols: ", "Municipal enforcement squads and tow-trucks patrol randomly, relying on guesswork rather than priority or traffic sensitivity data."),
        ("Resource Disconnection: ", "Highly accurate predictive traffic modeling exists, but is completely disconnected from real-time municipal dispatcher platforms.")
    ]
    add_textbox_bullets(slide2, Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.5), points2)

    # Image
    img2_path = r"docs/screenshots/urban_congestion_problem.png"
    if os.path.exists(img2_path):
        slide2.shapes.add_picture(img2_path, Inches(7.8), Inches(1.8), Inches(4.8), Inches(4.2))

    # -------------------------------------------------------------
    # SLIDE 3: The Solution
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide3)
    add_slide_header(slide3, "ParkFlow AI: End-to-End Predictive Dispatch")

    # 3 Column Cards
    col_width = Inches(3.6)
    gap = Inches(0.4)
    y_pos = Inches(2.0)
    card_h = Inches(4.2)

    # Card 1
    add_card(slide3, Inches(0.8), y_pos, col_width, card_h)
    tb1 = slide3.shapes.add_textbox(Inches(0.9), y_pos + Inches(0.2), col_width - Inches(0.2), card_h - Inches(0.4))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "1. Predictive ML Engine\n\n"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_CYAN
    r1 = p1.add_run()
    r1.text = "Combines LightGBM, XGBoost, and CatBoost ensemble. Dynamically forecasts spatiotemporal traffic demand for any geohash coordinate based on cyclic time parameters."
    r1.font.size = Pt(14)
    r1.font.color.rgb = TEXT_MUTED

    # Card 2
    add_card(slide3, Inches(0.8) + col_width + gap, y_pos, col_width, card_h)
    tb2 = slide3.shapes.add_textbox(Inches(0.9) + col_width + gap, y_pos + Inches(0.2), col_width - Inches(0.2), card_h - Inches(0.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "2. Bottleneck Engine\n\n"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_GREEN
    r2 = p2.add_run()
    r2.text = "Custom formulas evaluate congestion drag (CDI) caused by double-parked vehicles and compute dynamic targeted enforcement priority (TEP) scores for dispatch."
    r2.font.size = Pt(14)
    r2.font.color.rgb = TEXT_MUTED

    # Card 3
    add_card(slide3, Inches(0.8) + (col_width + gap) * 2, y_pos, col_width, card_h)
    tb3 = slide3.shapes.add_textbox(Inches(0.9) + (col_width + gap) * 2, y_pos + Inches(0.2), col_width - Inches(0.2), card_h - Inches(0.4))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "3. Actionable UI Dashboard\n\n"
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_CYAN
    r3 = p3.add_run()
    r3.text = "Interactive Leaflet mapping, auto-ranked priority dispatch list for enforcement officers, and an instant What-If congestion clearing simulator."
    r3.font.size = Pt(14)
    r3.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 4: ML Engine Deep-Dive (Phase 1)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide4)
    add_slide_header(slide4, "Phase 1 ML Architecture: Blended Ensemble")

    # Left box
    points4_left = [
        ("Blended Regressor Ensemble: ", "Built a powerful blend of LightGBM, XGBoost, and CatBoost. Blending optimized weights to achieve robust spatiotemporal generalizability."),
        ("Cyclic Temporal Features: ", "Encoded continuous time boundaries. Applied sine/cosine transformations on hour of day and day of week to preserve continuous temporal patterns.")
    ]
    add_textbox_bullets(slide4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5), points4_left)

    # Right box
    points4_right = [
        ("Spatial Feature Clustering: ", "Decoded 9-character spatial geohashes to precise latitude/longitude coordinates and grouped coordinates using K-Means clustering (K=15) to map regional densities."),
        ("Target Baseline Encoding: ", "Computed out-of-fold historical baseline demand per geohash/temporal bucket to provide baseline benchmarks without introducing data leakage.")
    ]
    add_textbox_bullets(slide4, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5), points4_right)

    # -------------------------------------------------------------
    # SLIDE 5: Formulas
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide5)
    add_slide_header(slide5, "Congestion Impact Formulas & Economic Loss")

    # 3 Columns for CDI, TEP, and Loss
    col_w5 = Inches(3.6)
    gap5 = Inches(0.4)
    y_pos5 = Inches(1.8)
    card_h5 = Inches(4.8)

    # CDI Box
    add_card(slide5, Inches(0.8), y_pos5, col_w5, card_h5)
    tb5_1 = slide5.shapes.add_textbox(Inches(0.9), y_pos5 + Inches(0.2), col_w5 - Inches(0.2), card_h5 - Inches(0.4))
    tf5_1 = tb5_1.text_frame
    tf5_1.word_wrap = True
    p5_1 = tf5_1.paragraphs[0]
    p5_1.text = "Congestion Drag Index\n(CDI)\n\n"
    p5_1.font.size = Pt(20)
    p5_1.font.bold = True
    p5_1.font.color.rgb = ACCENT_CYAN
    r5_1 = p5_1.add_run()
    r5_1.text = "CDI = predicted_demand / remaining_capacity\n\n* Where remaining_capacity decreases by 7% per double-parked vehicle (up to a max reduction of 65%).\n\n* Highlights road corridors that are approaching severe traffic saturation."
    r5_1.font.size = Pt(14)
    r5_1.font.color.rgb = TEXT_MUTED

    # TEP Box
    add_card(slide5, Inches(0.8) + col_w5 + gap5, y_pos5, col_w5, card_h5)
    tb5_2 = slide5.shapes.add_textbox(Inches(0.9) + col_w5 + gap5, y_pos5 + Inches(0.2), col_w5 - Inches(0.2), card_h5 - Inches(0.4))
    tf5_2 = tb5_2.text_frame
    tf5_2.word_wrap = True
    p5_2 = tf5_2.paragraphs[0]
    p5_2.text = "Enforcement Priority\n(TEP)\n\n"
    p5_2.font.size = Pt(20)
    p5_2.font.bold = True
    p5_2.font.color.rgb = ACCENT_GREEN
    r5_2 = p5_2.add_run()
    r5_2.text = "TEP = violation_count * CDI\n\n* Ranks active hotspots by both physical violations and the road's congestion index.\n\n* Enables dispatchers to target tow-trucks to corridors where they will relieve the most traffic drag."
    r5_2.font.size = Pt(14)
    r5_2.font.color.rgb = TEXT_MUTED

    # Loss Box
    add_card(slide5, Inches(0.8) + (col_w5 + gap5)*2, y_pos5, col_w5, card_h5)
    tb5_3 = slide5.shapes.add_textbox(Inches(0.9) + (col_w5 + gap5)*2, y_pos5 + Inches(0.2), col_w5 - Inches(0.2), card_h5 - Inches(0.4))
    tf5_3 = tb5_3.text_frame
    tf5_3.word_wrap = True
    p5_3 = tf5_3.paragraphs[0]
    p5_3.text = "Economic Loss\n(Rupees)\n\n"
    p5_3.font.size = Pt(20)
    p5_3.font.bold = True
    p5_3.font.color.rgb = ACCENT_CYAN
    r5_3 = p5_3.add_run()
    r5_3.text = "Loss = delay_minutes * affected_vehicles * hourly_rate\n\n* Estimates financial commuter drain.\n\n* Assumes ₹300 per hour commuter time-value and 400 affected vehicles per hour. Translates physical traffic blockages into clear, understandable economic figures (₹)."
    r5_3.font.size = Pt(14)
    r5_3.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 6: Interactive Dashboard
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide6)
    add_slide_header(slide6, "Operational Dashboard: Map & Priority Dispatch")

    # Left bullets
    points6 = [
        ("Dynamic Leaflet Mapping: ", "Renders coordinates in Bengaluru. Color-coded markers (Red = High Congestion, Yellow = Moderate, Green = Low) reveal spatial parking hotspot clusters."),
        ("Auto-Ranked Dispatch: ", "Left sidebar lists hotspots in descending order of their TEP score, showing officers exactly where to deploy for maximum congestion relief."),
        ("Diurnal Analytics Chart: ", "Displays city-wide traffic demand patterns to schedule patrol rosters efficiently.")
    ]
    add_textbox_bullets(slide6, Inches(0.8), Inches(1.8), Inches(6.2), Inches(4.5), points6)

    # Dashboard screenshot
    img6_path = r"docs/screenshots/dashboard_overview.png"
    if os.path.exists(img6_path):
        slide6.shapes.add_picture(img6_path, Inches(7.4), Inches(1.8), Inches(5.1), Inches(4.2))

    # -------------------------------------------------------------
    # SLIDE 7: What-If Simulator
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide7)
    add_slide_header(slide7, "What-If Congestion Simulator")

    # Left bullets
    points7 = [
        ("Simulate Vehicle Clearing: ", "Allows operators to slide from 0 to maximum parked vehicles to simulate active traffic enforcement actions."),
        ("Instant Financial Feedback: ", "Recalculates in real-time the capacity restored, delay minutes saved, and the total commuter economic savings (₹) achieved."),
        ("Corridor Demand Analytics: ", "Includes detailed Chart.js visual showing the selected location's violation history alongside the city average demand curve.")
    ]
    add_textbox_bullets(slide7, Inches(0.8), Inches(1.8), Inches(6.2), Inches(4.5), points7)

    # Simulator screenshot
    img7_path = r"docs/screenshots/location_inspector.png"
    if os.path.exists(img7_path):
        slide7.shapes.add_picture(img7_path, Inches(7.4), Inches(1.8), Inches(5.1), Inches(4.2))

    # -------------------------------------------------------------
    # SLIDE 8: Value & Future Roadmap
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide8)
    add_slide_header(slide8, "Business Value & Future Scaling Roadmap")

    # Left Column
    points8_left = [
        ("Commuter Travel Savings: ", "Reduces travel delays, carbon emissions, and fuel burn in dense urban bottlenecks by maximizing road capacity."),
        ("Enforcement Cost Relief: ", "Allows municipalities to maximize tow-truck ROI. deploys squads to high-impact priority grids instead of blind random patrols.")
    ]
    add_textbox_bullets(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5), points8_left)

    # Right Column
    points8_right = [
        ("IoT & Smart City Sensors: ", "Plan to integrate low-cost geomagnetic road sensors that report occupancy, automating violation logging without manual entry."),
        ("Edge-AI CCTV Camera Feeds: ", "Roadmap to integrate virtual no-parking boxes on city CCTV cameras, running YOLOv8 model for real-time automated vehicle counts.")
    ]
    add_textbox_bullets(slide8, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5), points8_right)

    # Save
    out_path = r"ParkFlow_AI_Presentation.pptx"
    prs.save(out_path)
    print(f"Presentation generated successfully at: {os.path.abspath(out_path)}")

if __name__ == "__main__":
    create_presentation()
